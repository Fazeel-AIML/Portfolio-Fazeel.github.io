import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

def add_image_with_aspect_ratio(slide, img_path, left, top, max_width, max_height):
    try:
        with Image.open(img_path) as img:
            orig_width, orig_height = img.size
        
        aspect_ratio = orig_width / orig_height
        target_aspect = max_width.inches / max_height.inches

        if aspect_ratio > target_aspect:
            final_width = max_width.inches
            final_height = final_width / aspect_ratio
        else:
            final_height = max_height.inches
            final_width = final_height * aspect_ratio

        img_left = left.inches + (max_width.inches - final_width) / 2
        img_top = top.inches + (max_height.inches - final_height) / 2

        slide.shapes.add_picture(img_path, Inches(img_left), Inches(img_top), width=Inches(final_width), height=Inches(final_height))
    except Exception as e:
        print(f"Error adding {img_path}: {e}")

def add_tech_tags(slide, tags, start_left, start_top, spacing=0.2):
    current_top = start_top
    for tag in tags:
        # Approximate dynamic width relative to text length
        approx_width = Inches(max(1.3, len(tag) * 0.12))
        approx_height = Inches(0.4)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            start_left, 
            current_top, 
            approx_width, 
            approx_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(14, 165, 233) # Professional sky/blue highlight
        shape.line.color.rgb = RGBColor(14, 165, 233)
        
        tf = shape.text_frame
        tf.word_wrap = False
        
        # Center horizontally & vertically inside the rectangle
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        current_top += approx_height + Inches(spacing)

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42) # Deep Slate/Navy background
    bg.line.color.rgb = RGBColor(15, 23, 42)

    # Watermark
    wm_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(3.5))
    wm_box.rotation = -20
    p_wm = wm_box.text_frame.add_paragraph()
    p_wm.text = "7 Kings Code"
    p_wm.font.bold = True
    p_wm.font.size = Pt(110)
    p_wm.font.color.rgb = RGBColor(25, 33, 52) # extremely faint dark blue for watermark
    p_wm.alignment = PP_ALIGN.CENTER
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "AI & Engineering Portfolio"
    p.font.bold = True
    p.font.size = Pt(56)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "A Collection of Production-Ready Systems"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.add_paragraph()
    p_sub.text = "7 Kings Code"
    p_sub.font.bold = True
    p_sub.font.size = Pt(30)
    p_sub.font.color.rgb = RGBColor(56, 189, 248) # Sky blue accent
    p_sub.alignment = PP_ALIGN.CENTER

    projects = [
        {
            "title": "Autonomous Car Simulation",
            "desc": "Real-time autonomous navigation system integrating obstacle detection and GPS tracking. Trained on custom simulated driving datasets annotated via Roboflow. Utilizes Convolutional Neural Networks (CNNs) for lane tracking and YOLO object detection models for dynamic obstacle avoidance. Controlled via Django edge computing logic.",
            "img": "assets/img/auto_car - Copy.GIF",
            "tech": ["Raspberry Pi", "OpenCV", "Django", "Machine Learning"]
        },
        {
            "title": "SignoraAI - Sign Recognition",
            "desc": "An advanced computer vision system dedicated to sign language recognition. Trained on custom gesture datasets meticulously annotated using Roboflow. Utilizes state-of-the-art YOLOv8 deep learning models alongside MediaPipe for precise real-time hand, facial, and pose tracking classification.",
            "img": "assets/img/Signora.gif",
            "tech": ["Computer Vision", "Deep Learning", "CNN"]
        },
        {
            "title": "Medical Chatbot with Voice",
            "desc": "Next-generation voice-enabled medical consultation chatbot. Leveraging RAG (Retrieval-Augmented Generation) architecture utilizing vector databases storing thousands of medical journals. Uses OpenAI's Whisper for speech-to-text processing and fine-tuned LLMs for accurate health advice delivery.",
            "img": "assets/img/Medical_Chatbot2.gif",
            "tech": ["LangChain", "Voice AI", "RAG"]
        },
        {
            "title": "Email Classifier & Ticket Generator",
            "desc": "Fully autonomous agentic system built to streamline IT workflows. Trained on diverse datasets of internal IT support emails. Employs Natural Language Processing algorithms (BERT/RoBERTa) for intent recognition, urgency classification, and automatic entity extraction to generate organized support tickets.",
            "img": "assets/img/Tickek.gif",
            "tech": ["Agentic AI", "NLP", "Automation"]
        },
        {
            "title": "Aya - Emotion Detection ChatBot",
            "desc": "An emotionally intelligent conversational agent trained on vast sentiment analysis datasets such as GoEmotions. Incorporates advanced Transformer-based NLP architectures to parse nuanced user phrasing, accurately gauge emotional states, and dynamically construct highly empathetic responses.",
            "img": "assets/img/AYA.gif",
            "tech": ["Sentiment Analysis", "NLP", "LLM"]
        },
        {
            "title": "Live RAG Blogs ChatBot",
            "desc": "A dynamic document-querying engine utilizing Retrieval-Augmented Generation. Actively scrapes live blog content, mapping textual data into high-dimensional embeddings using HuggingFace sentence-transformers. Operates via vector databases like Pinecone to enable instant semantic searches.",
            "img": "assets/img/Blogs_Chatbot.gif",
            "tech": ["RAG", "LangChain", "Vector Databases"]
        },
        {
            "title": "Fine-Tuned Microsoft PHI (Medical)",
            "desc": "Deeply fine-tuned instance of the Microsoft PHI foundational model on specialized datasets such as PubMedQA. Employs sophisticated Parameter-Efficient Fine-Tuning (PEFT) and LoRA techniques to deliver confident and analytically precise clinical answers under severe computational constraints.",
            "img": "assets/img/Medical_Chatbot.gif",
            "tech": ["Fine-Tuning", "Microsoft PHI", "Healthcare AI"]
        },
        {
            "title": "Tumor Detection Model",
            "desc": "End-to-end brain anomaly classification pipeline using deep transfer learning protocols (ResNet50). Trained on challenging neurological imaging datasets (e.g., BraTS). Tailor-made CNN architectures intertwined with image segmentation methodologies (U-Net) swiftly pinpoint regions of interest within heavy MRI scans.",
            "img": "assets/img/Tumor_Detection.gif",
            "tech": ["Transfer Learning", "CNN", "Medical Imaging"]
        },
        {
            "title": "Vehicle Insurance Prediction System",
            "desc": "Highly scalable machine learning platform purposed for forecasting vehicle insurance claims. Developed using comprehensive tabular datasets representing historical claim data. Utilizes ensemble tree-based models like XGBoost/LightGBM. Equipped with CI/CD automation and AWS Docker deployment.",
            "img": "assets/img/veh_ins - Copy.GIF",
            "tech": ["MLOps", "Docker", "AWS", "CI/CD"]
        },
        {
            "title": "Weather Forecasting MLOps Pipeline",
            "desc": "Enterprise-grade MLOps deployment devoted to proactive weather predictions. Trained on historical meteorological time-series datasets. Utilizes Recurrent Neural Networks (LSTM) and Prophet anomaly forecasting techniques. Features automated retraining cycles with resilient cloud inference on AWS.",
            "img": "assets/img/weather_forcast.GIF",
            "tech": ["MLOps", "Cloud APIs", "AWS"]
        },
        {
            "title": "Fine-Tuned GPT Large Language Model",
            "desc": "Context-aware LLM calibrated meticulously for structured university ecosystems. Fine-tuned on custom datasets comprised of massive university course catalogs and institutional FAQs. Leverages targeted instruction-tuning techniques resulting in natural dialogue execution across complex academic queries.",
            "img": "assets/img/llm_uni - Copy.GIF",
            "tech": ["Fine-Tuning", "GPT", "Education Tech"]
        },
        {
            "title": "Escalator Traffic Surveillance",
            "desc": "Rapid computer-vision traffic dashboard observing real-time mall densities. Trained meticulously on a custom surveillance dataset annotated precisely via Roboflow. Fueled by customized YOLOv8 object detection deployments paired with DeepSORT tracking topologies to minimize spatial occlusion limits.",
            "img": "assets/img/esclatorscount.GIF",
            "tech": ["YOLOv8", "Object Detection", "Smart Cities"]
        },
        {
            "title": "Highway Vehicle Monitoring Array",
            "desc": "Heavy-duty tracking network running on robust urban and freeway video datasets prepared through Roboflow annotations. Employs lightweight YOLOv8n object detection checkpoints merged with ByteTrack algorithms to record car densities accurately despite extreme multi-target saturation.",
            "img": "assets/img/CarCounter.gif",
            "tech": ["YOLOv8n", "Tracking Algorithms", "Analytics"]
        },
        {
            "title": "Pathological Skin Condition Classifier",
            "desc": "Augmented dermatological reasoning model developed using substantial skin pathology datasets (e.g., HAM10000). Utilizes advanced progressive neural architectures (EfficientNet/VGG16) alongside rigorous data augmentation techniques for categorizing complex lesions into comprehensive clinical subclasses.",
            "img": "assets/img/skin_disease.gif",
            "tech": ["CNN Deep Learning", "Dermatology AI", "Computer Vision"]
        },
        {
            "title": "Urinary Sediment Particle Extraction",
            "desc": "Next-generation microscopic pattern recognizer designed for fluid analysis. Trained on custom high-resolution microscopic imagery data mapped and segmented utilizing Roboflow. Enhances laboratories utilizing YOLOv8 instance segmentation methodologies to rapidly extract sedimentary details under harsh visualization.",
            "img": "assets/img/Selected-samples-of-urinary-sediment-particle.png",
            "tech": ["Microscopy AI", "Deep Learning", "Biotech Systems"]
        }
    ]

    base_dir = r"d:\Fazeels_WorkSpace\Portfolio-Fazeel.github.io"
    
    # --- Content slides ---
    for proj in projects:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        
        # Background Watermark
        wm_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(3.5))
        wm_box.rotation = -20
        p_wm = wm_box.text_frame.add_paragraph()
        p_wm.text = "7 Kings Code"
        p_wm.font.bold = True
        p_wm.font.size = Pt(110)
        p_wm.font.color.rgb = RGBColor(245, 245, 245) # very faint grey
        p_wm.alignment = PP_ALIGN.CENTER
        
        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = proj["title"]
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(15, 23, 42)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(1.8))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.add_paragraph()
        p_desc.text = proj["desc"]
        p_desc.font.size = Pt(17)
        p_desc.font.color.rgb = RGBColor(71, 85, 105)
        
        # Divider Line
        line = slide.shapes.add_shape(
            9, Inches(0.5), Inches(3.2), Inches(12.333), Inches(0) # 9 is MSO_SHAPE.LINE
        )
        line.line.color.rgb = RGBColor(226, 232, 240)
        
        # Render dynamic visual tags (Pills vertically aligned on the left)
        add_tech_tags(slide, proj["tech"], Inches(0.5), Inches(3.5), spacing=0.3)
        
        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.5))
        tf_footer = footer_box.text_frame
        p_footer = tf_footer.add_paragraph()
        p_footer.text = "7 Kings Code"
        p_footer.font.bold = True
        p_footer.font.size = Pt(14)
        p_footer.font.color.rgb = RGBColor(200, 200, 200)
        p_footer.alignment = PP_ALIGN.RIGHT
        
        # Add image/gif bounded right of the tags
        img_path = os.path.join(base_dir, proj["img"])
        if os.path.exists(img_path):
            add_image_with_aspect_ratio(
                slide, 
                img_path, 
                left=Inches(2.5), 
                top=Inches(3.3),   
                max_width=Inches(10.3), 
                max_height=Inches(4.1)
            )
        else:
            print(f"File not found: {img_path}")
            
    # Save the PPTX
    output_path = os.path.join(base_dir, "Deep_Learning_Projects_7KingsCode_V4.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == '__main__':
    build_presentation()
